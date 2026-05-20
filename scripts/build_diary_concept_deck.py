#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_diary_concept_deck.py
============================
인생포트폴리오 맞춤형 다이어리 — 시안 PPT 빌드 스크립트 (16 슬라이드, v1.3)

입력: HANDOFF_v2.md + 02_concept_deck_outline.md (사양 참조)
출력: docs/strategy/diary/manufacturer-brief/02_concept_deck.pptx

설계 원칙:
- Apple Keynote 톤. 미니멀 텍스트 중심. 흰 배경 + 네이비/골드.
- 슬라이드별로 페이지가 잘리지 않게 명확히 1슬라이드 = 1주제.
- 한글 폰트: Noto Sans CJK KR / Malgun Gothic 폴백.
- 브랜드 컬러: 짙은 네이비 #1A2B4A, 골드 #C9A04F.
- 노출 범위 원칙: 가격/마진/매출 일체 미노출.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ============================================================
# 상수
# ============================================================
ROOT = Path("/home/user/webapp")
OUTPUT_PPTX = ROOT / "docs/strategy/diary/manufacturer-brief/02_concept_deck.pptx"

# 슬라이드 크기 (16:9 와이드)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 컬러
NAVY = RGBColor(0x1A, 0x2B, 0x4A)
GOLD = RGBColor(0xC9, 0xA0, 0x4F)
GRAY_900 = RGBColor(0x1F, 0x29, 0x37)
GRAY_700 = RGBColor(0x37, 0x41, 0x51)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
GRAY_400 = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_200 = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_50  = RGBColor(0xF9, 0xFA, 0xFB)
CREAM    = RGBColor(0xFE, 0xF7, 0xE6)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# 폰트
FONT_KR = "Noto Sans CJK KR"
FONT_KR_FALLBACK = "Malgun Gothic"


# ============================================================
# 헬퍼 함수
# ============================================================
def set_font(run, *, size=14, bold=False, color=None, name=FONT_KR):
    """폰트 일괄 설정."""
    run.font.name = name
    # 동아시아 글꼴(eastAsia)도 같이 지정 — 한글 안정 렌더링
    rPr = run._r.get_or_add_rPr()
    # 기존 eastAsia 제거
    for ea in rPr.findall(qn("a:ea")):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", name)
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text="",
                size=14, bold=False, color=GRAY_900,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                line_spacing=1.3):
    """텍스트박스 추가 (단일 paragraph)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return tb


def add_multitext(slide, left, top, width, height, lines,
                  anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
                  line_spacing=1.4):
    """
    여러 단락 텍스트박스.
    lines: [(text, size, bold, color), ...]
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    for i, item in enumerate(lines):
        text, size, bold, color = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        set_font(run, size=size, bold=bold, color=color)
    return tb


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=0.75):
    """직사각형(라인/필) 추가."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    return shape


def add_footer(slide, page_num, total=16, show_meta=True):
    """모든 슬라이드 공통 푸터."""
    # 푸터 위 가는 라인
    line = add_rect(slide, Inches(0.6), Inches(7.05), Inches(12.133), Emu(4500),
                    fill_color=GRAY_200, line_color=None)

    # 좌측: 브랜드
    if show_meta:
        add_textbox(slide, Inches(0.6), Inches(7.15), Inches(7.0), Inches(0.3),
                    text="Life Portfolio · 인생포트폴리오 맞춤형 다이어리",
                    size=8.5, color=GRAY_500)

        # 우측: 버전 + 페이지
        add_textbox(slide, Inches(8.0), Inches(7.15), Inches(4.7), Inches(0.3),
                    text=f"v1.3 · 500부 견적 요청    |    {page_num} / {total}",
                    size=8.5, color=GRAY_500, align=PP_ALIGN.RIGHT)


def add_slide_header(slide, slide_num, kicker, title):
    """페이지 상단 헤더(키커 + 큰 타이틀)."""
    # 좌측 골드 바
    add_rect(slide, Inches(0.6), Inches(0.7), Inches(0.08), Inches(0.45),
             fill_color=GOLD, line_color=None)

    # 키커(섹션명)
    add_textbox(slide, Inches(0.85), Inches(0.68), Inches(8.0), Inches(0.25),
                text=kicker.upper(), size=9, bold=True, color=GOLD)

    # 타이틀
    add_textbox(slide, Inches(0.85), Inches(0.92), Inches(11.5), Inches(0.6),
                text=title, size=24, bold=True, color=NAVY)

    # 우측 페이지 번호 (대형)
    add_textbox(slide, Inches(11.3), Inches(0.62), Inches(1.4), Inches(0.55),
                text=f"{slide_num:02d}", size=28, bold=True, color=GRAY_200,
                align=PP_ALIGN.RIGHT)

    # 타이틀 하단 굵은 라인
    add_rect(slide, Inches(0.6), Inches(1.65), Inches(2.0), Inches(0.03),
             fill_color=NAVY, line_color=None)
    add_rect(slide, Inches(2.6), Inches(1.65), Inches(10.13), Emu(7000),
             fill_color=GRAY_200, line_color=None)


def add_blank_slide(prs):
    """완전 빈 슬라이드(레이아웃 6번 = blank) 추가."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    return slide


# ============================================================
# 슬라이드 1: 표지
# ============================================================
def slide_01_cover(prs):
    slide = add_blank_slide(prs)

    # v1.3: 풀 네이비 배경 (실제 다이어리 표지처럼)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H,
             fill_color=NAVY, line_color=None)

    # 상하 골드 라인 (얇은 프레임)
    add_rect(slide, Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.025),
             fill_color=GOLD, line_color=None)
    add_rect(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.025),
             fill_color=GOLD, line_color=None)

    # v1.3: L 모노그램 박스 제거 (고객 관점에서 불필요)
    # 상단 골드 디바이더 (브랜드 텍스트 위)
    add_rect(slide, Inches(6.4), Inches(2.9), Inches(0.5), Inches(0.015),
             fill_color=GOLD, line_color=None)

    # 영문 브랜드 (중앙으로 재배치)
    add_textbox(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.5),
                text="LIFE · PORTFOLIO",
                size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # 한글 브랜드 (메인 타이틀)
    add_textbox(slide, Inches(0.8), Inches(3.95), Inches(11.5), Inches(1.0),
                text="인생포트폴리오",
                size=54, bold=True, color=CREAM, align=PP_ALIGN.CENTER)

    # 골드 디바이더 (중앙)
    add_rect(slide, Inches(6.15), Inches(5.35), Inches(1.0), Inches(0.025),
             fill_color=GOLD, line_color=None)

    # 'Only One' — 핵심 차별점 (v1.2 신규, v1.3 유지)
    add_textbox(slide, Inches(0.8), Inches(5.55), Inches(11.5), Inches(0.6),
                text="Only One",
                size=26, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

    # v1.3: 표지에는 페이지 번호 / 연도 / 메타 / 종이 / 사이즈 / L 로고 일체 미노출
    # 사용자 요청: 고객 관점에서 불필요한 정보 일체 제거
    return slide


# ============================================================
# 슬라이드 2: 목차
# ============================================================
def slide_02_toc(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 2, "Contents", "목차 — 5개 챕터")

    chapters = [
        ("01", "제품 개요", "한 줄 정의 · 타겟 사용자 · 사용 시나리오", "3 ~ 4"),
        ("02", "외관 사양", "사이즈 · 표지 · 컬러 · 종이 · 마감", "5 ~ 7"),
        ("03", "속지 구조", "Part 0 ~ Part 7 · 합계 256p · 연간/월간/메모장 포함", "8 ~ 13"),
        ("04", "인쇄 구조", "일반 인쇄 단일 층 / 손글씨 친화 타이포", "14 ~ 15"),
        ("05", "견적 요청", "8개 항목 + 회신 양식 + 일정", "16"),
    ]

    top = Inches(2.1)
    row_h = Inches(0.85)
    for i, (num, title, desc, pages) in enumerate(chapters):
        y = top + row_h * i

        # 챕터 번호
        add_textbox(slide, Inches(0.85), y, Inches(0.9), Inches(0.5),
                    text=num, size=28, bold=True, color=GOLD)
        # 챕터 제목
        add_textbox(slide, Inches(2.0), y + Inches(0.05), Inches(4.5), Inches(0.45),
                    text=title, size=18, bold=True, color=NAVY)
        # 설명
        add_textbox(slide, Inches(6.5), y + Inches(0.1), Inches(4.5), Inches(0.4),
                    text=desc, size=11, color=GRAY_500)
        # 슬라이드 범위
        add_textbox(slide, Inches(11.2), y + Inches(0.1), Inches(1.5), Inches(0.4),
                    text=pages, size=11, color=GRAY_400, align=PP_ALIGN.RIGHT)
        # 구분선
        if i < len(chapters) - 1:
            add_rect(slide, Inches(0.85), y + Inches(0.7),
                     Inches(11.85), Emu(4000),
                     fill_color=GRAY_200, line_color=None)

    add_footer(slide, 2)


# ============================================================
# 슬라이드 3: 제품 한 줄 정의
# ============================================================
def slide_03_definition(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 3, "01 · 제품 개요", "제품 한 줄 정의")

    # 큰 따옴표 블록
    add_textbox(slide, Inches(0.85), Inches(2.2), Inches(0.8), Inches(1.5),
                text="“", size=90, bold=True, color=GOLD)

    # 핵심 문장 (크게)
    add_textbox(slide, Inches(1.8), Inches(2.5), Inches(11.0), Inches(2.2),
                text="자기 검사 리포트가 그대로\n옮겨 적히는 1년치 자기설계 다이어리.",
                size=32, bold=True, color=NAVY, line_spacing=1.35)

    # 골드 라인
    add_rect(slide, Inches(1.8), Inches(5.2), Inches(2.5), Inches(0.04),
             fill_color=GOLD, line_color=None)

    # 보조 설명
    add_textbox(slide, Inches(1.8), Inches(5.4), Inches(10.5), Inches(1.5),
                text="검사 리포트의 4SE(자기이해/자기표현/자기설계/자기실행)와 실행 프로그램의 3주·3개월·1년 구조를 그대로 받아 적는 빈칸 양식을 갖춘다. 다이어리는 리포트의 연장선이며, 두 산출물의 라벨·항목명·순서는 100% 일치한다.",
                size=13, color=GRAY_700, line_spacing=1.7)

    add_footer(slide, 3)


# ============================================================
# 슬라이드 4: 타겟 사용자 & 사용 시나리오
# ============================================================
def slide_04_target(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 4, "01 · 제품 개요", "타겟 사용자 & 사용 시나리오")

    # 왼쪽: 타겟 사용자
    add_textbox(slide, Inches(0.85), Inches(2.1), Inches(5.8), Inches(0.4),
                text="타겟 사용자", size=14, bold=True, color=NAVY)
    add_rect(slide, Inches(0.85), Inches(2.5), Inches(1.0), Inches(0.03),
             fill_color=GOLD, line_color=None)

    target_lines = [
        ("· 30 ~ 50대 자기설계 의지자", 13, False, GRAY_900),
        ("· 검사·진단 도구를 일상에 통합하려는 직장인 / 사업자", 13, False, GRAY_900),
        ("· 1년 단위 목표 설계 + 주간 실행 점검을 병행하려는 사용자", 13, False, GRAY_900),
        ("· 만년필·종이 다이어리를 선호하는 아날로그 친화 사용자", 13, False, GRAY_900),
        ("· 신앙/철학/인문적 분기점을 일상에 두려는 독자", 13, False, GRAY_900),
    ]
    add_multitext(slide, Inches(0.85), Inches(2.75), Inches(5.8), Inches(3.5),
                  lines=target_lines, line_spacing=1.85)

    # 오른쪽: 사용 시나리오 (3 step)
    add_textbox(slide, Inches(7.2), Inches(2.1), Inches(5.5), Inches(0.4),
                text="사용 시나리오 (3 step)", size=14, bold=True, color=NAVY)
    add_rect(slide, Inches(7.2), Inches(2.5), Inches(1.0), Inches(0.03),
             fill_color=GOLD, line_color=None)

    steps = [
        ("STEP 1", "검사 응답", "온라인 검사 도구로 본인의 사명·비전·4SE 결과 도출"),
        ("STEP 2", "리포트 수령", "PDF 검사 리포트 5p + 실행 프로그램 PDF 6p 발급"),
        ("STEP 3", "다이어리 작성", "다이어리 Part 0~7의 빈칸 양식에 리포트 결과를 옮겨 적음"),
    ]
    y0 = Inches(2.8)
    for i, (st, label, desc) in enumerate(steps):
        y = y0 + Inches(1.05) * i
        # 단계 번호 박스
        add_rect(slide, Inches(7.2), y, Inches(0.9), Inches(0.85),
                 fill_color=NAVY, line_color=None)
        add_textbox(slide, Inches(7.2), y + Inches(0.18), Inches(0.9), Inches(0.5),
                    text=st, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # 라벨
        add_textbox(slide, Inches(8.3), y + Inches(0.05), Inches(4.5), Inches(0.35),
                    text=label, size=13, bold=True, color=NAVY)
        # 설명
        add_textbox(slide, Inches(8.3), y + Inches(0.4), Inches(4.5), Inches(0.5),
                    text=desc, size=10.5, color=GRAY_700, line_spacing=1.4)

    add_footer(slide, 4)


# ============================================================
# 슬라이드 5: 외관 사양
# ============================================================
def slide_05_appearance(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 5, "02 · 외관 사양", "외관 사양 한눈에 보기")

    items = [
        ("사이즈", "A5 (148 × 210mm)", "몰스킨 위클리 라지와 동일 폭"),
        ("표지", "PU 양장 + 금박 박표지", "형압 인생포트폴리오 로고"),
        ("제본", "양장 사철 (Smyth-sewn)", "180° 펼침 — 만년필 작성 친화"),
        ("부속", "가름끈 2개 / 모서리 코너 4개", "인덱스 박 / 뒷주머니 1개"),
        ("포장", "박스 + 인삿말 카드", "+ 사용 가이드 8p"),
        ("페이지", "256p (16의 배수 · 16절판 최적)", "속표지·간지 8p + 본문 248p — 70g 단일 견적"),
    ]

    # 2 × 3 그리드
    col_w = Inches(4.1)
    row_h = Inches(1.55)
    x0 = Inches(0.6)
    y0 = Inches(2.1)
    gap_x = Inches(0.15)
    gap_y = Inches(0.18)

    for i, (label, value, desc) in enumerate(items):
        col = i % 3
        row = i // 3
        x = x0 + (col_w + gap_x) * col
        y = y0 + (row_h + gap_y) * row

        # 카드 박스
        add_rect(slide, x, y, col_w, row_h,
                 fill_color=GRAY_50, line_color=GRAY_200, line_width=0.5)
        # 좌측 골드 바
        add_rect(slide, x, y, Inches(0.06), row_h,
                 fill_color=GOLD, line_color=None)

        # 라벨
        add_textbox(slide, x + Inches(0.25), y + Inches(0.15),
                    col_w - Inches(0.4), Inches(0.3),
                    text=label, size=10, bold=True, color=GOLD)
        # 값
        add_textbox(slide, x + Inches(0.25), y + Inches(0.45),
                    col_w - Inches(0.4), Inches(0.45),
                    text=value, size=14, bold=True, color=NAVY)
        # 설명
        add_textbox(slide, x + Inches(0.25), y + Inches(0.95),
                    col_w - Inches(0.4), Inches(0.5),
                    text=desc, size=10, color=GRAY_500, line_spacing=1.4)

    add_footer(slide, 5)


# ============================================================
# 슬라이드 6: 컬러 & 마감
# ============================================================
def slide_06_color(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 6, "02 · 외관 사양", "컬러 & 마감 옵션")

    # 4가지 컬러 칩
    colors = [
        ("CLASSIC NAVY", "클래식 네이비", "기본 컬러", RGBColor(0x1A, 0x2B, 0x4A)),
        ("BURGUNDY",     "버건디",       "옵션 1",    RGBColor(0x6B, 0x1F, 0x2A)),
        ("FOREST GREEN", "그린",         "옵션 2",    RGBColor(0x1F, 0x3F, 0x32)),
        ("LIMITED EDITION", "한정판",     "별도 협의", RGBColor(0x3D, 0x2E, 0x1F)),
    ]

    chip_w = Inches(2.85)
    chip_h = Inches(3.2)
    x0 = Inches(0.6)
    y0 = Inches(2.1)
    gap = Inches(0.2)

    for i, (en, kr, tag, rgb) in enumerate(colors):
        x = x0 + (chip_w + gap) * i

        # 컬러 칩(큰 사각형)
        add_rect(slide, x, y0, chip_w, Inches(2.1),
                 fill_color=rgb, line_color=None)

        # 컬러 라벨 박스(아래 흰 영역)
        add_rect(slide, x, y0 + Inches(2.1), chip_w, Inches(1.1),
                 fill_color=WHITE, line_color=GRAY_200, line_width=0.5)

        # 영문명
        add_textbox(slide, x + Inches(0.2), y0 + Inches(2.2),
                    chip_w - Inches(0.4), Inches(0.3),
                    text=en, size=9, bold=True, color=GOLD)
        # 한글명
        add_textbox(slide, x + Inches(0.2), y0 + Inches(2.5),
                    chip_w - Inches(0.4), Inches(0.4),
                    text=kr, size=15, bold=True, color=NAVY)
        # 태그
        add_textbox(slide, x + Inches(0.2), y0 + Inches(2.9),
                    chip_w - Inches(0.4), Inches(0.3),
                    text=tag, size=10, color=GRAY_500)

    # 하단 마감 설명
    add_rect(slide, Inches(0.6), Inches(5.6), Inches(12.13), Inches(1.3),
             fill_color=CREAM, line_color=None)
    add_rect(slide, Inches(0.6), Inches(5.6), Inches(0.06), Inches(1.3),
             fill_color=GOLD, line_color=None)

    finish_text = [
        ("표지 마감 사양", 11, True, GOLD),
        ("· 표지 본체: PU 양장 (제조사 추천 자재 1종 + 대안 1종 제시 요청)", 11, False, GRAY_900),
        ("· 박표지: 골드 박(#C9A04F 톤) + 형압(blind emboss) 인생포트폴리오 로고", 11, False, GRAY_900),
        ("· 측면(에지): 무가공 — 만년필 친화 + 페이지 펼침 안정성 우선", 11, False, GRAY_900),
    ]
    add_multitext(slide, Inches(0.85), Inches(5.75), Inches(11.7), Inches(1.1),
                  lines=finish_text, line_spacing=1.5)

    add_footer(slide, 6)


# ============================================================
# 슬라이드 7: 종이 사양 양자 견적
# ============================================================
def slide_07_paper(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 7, "02 · 외관 사양", "종이 사양 — 만년필 친화 70g 단일 견적")

    # 핵심 메시지
    add_textbox(slide, Inches(0.85), Inches(2.05), Inches(12.0), Inches(0.5),
                text="손글씨 다이어리 — 만년필·중성펜 친화 70g 단일 시나리오 (v1.3 확정)",
                size=12, bold=True, color=GRAY_700)

    # 메인 카드 — 70g 단일
    card_x = Inches(0.6)
    card_y = Inches(2.7)
    card_w = Inches(12.13)
    card_h = Inches(4.3)
    add_rect(slide, card_x, card_y, card_w, card_h,
             fill_color=WHITE, line_color=GRAY_200, line_width=1.0)
    # 상단 골드 바
    add_rect(slide, card_x, card_y, card_w, Inches(0.7),
             fill_color=GOLD, line_color=None)
    add_textbox(slide, card_x + Inches(0.3), card_y + Inches(0.16), card_w - Inches(0.6), Inches(0.5),
                text="만년필 친화 70g — 도모에리버 / OK Fools 대체재",
                size=16, bold=True, color=NAVY)

    # 좌측 — 규격
    paper_spec = [
        ("규격", 11, True, GOLD),
        ("· 평량: 70gsm", 11, False, GRAY_900),
        ("· 색상: 아이보리 / 미색", 11, False, GRAY_900),
        ("· 두께감: 적정 — 256p 기준 책등 약 20mm", 11, False, GRAY_900),
        ("· 비침: 거의 없음 — 양면 필기 안정", 11, False, GRAY_900),
        ("· 만년필 잉크 번짐: 최소", 11, False, GRAY_900),
    ]
    add_multitext(slide, card_x + Inches(0.4), card_y + Inches(0.95),
                  Inches(5.5), Inches(3.2),
                  lines=paper_spec, line_spacing=1.65)

    # 우측 — 선택 이유
    paper_reason = [
        ("선택 이유 (v1.3)", 11, True, GOLD),
        ("· Part 0 7항목·Part 3 자유 메모를 손글씨로", 11, False, GRAY_900),
        ("  채우는 \"손글씨 다이어리\" 컨셉 — 70g 단일 최적", 11, False, GRAY_900),
        ("· 80g 보급형은 만년필 사용 시 비침 발생 → 제외", 11, False, GRAY_900),
        ("· 256p × 70g = 휴대성 + 1년 기록 충분 (몰스킨급)", 11, False, GRAY_900),
        ("· 견적 단일화로 발주 의사결정 단축", 11, True, NAVY),
    ]
    add_multitext(slide, card_x + Inches(6.2), card_y + Inches(0.95),
                  Inches(5.5), Inches(3.2),
                  lines=paper_reason, line_spacing=1.65)

    add_footer(slide, 7)


# ============================================================
# 슬라이드 8: 속지 구조 개요
# ============================================================
def slide_08_inner_overview(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 8, "03 · 속지 구조", "속지 구조 개요 — Part 0 ~ Part 7")

    # 좌측: 텍스트 가이드
    add_textbox(slide, Inches(0.85), Inches(2.05), Inches(5.8), Inches(0.4),
                text="구조 원칙", size=12, bold=True, color=NAVY)
    add_rect(slide, Inches(0.85), Inches(2.45), Inches(0.8), Inches(0.03),
             fill_color=GOLD, line_color=None)

    principle_lines = [
        ("· 검사 리포트 5p + 실행 프로그램 6p의 항목명을\n  다이어리 빈칸 양식에 그대로 반영", 11, False, GRAY_900),
        ("· 4SE(자기이해·자기표현·자기설계·자기실행)와\n  4라벨(본질·한 마디·키워드·결합) 공통 유지", 11, False, GRAY_900),
        ("· 페이지 합계 256p (16의 배수 · 16절판)\n  속표지·간지 8p + 본문 248p", 11, False, GRAY_900),
        ("· 전 부 일반 인쇄 — VDP 없음\n  Part 0 7항목은 사용자가 손글씨로 옮겨 적음", 11, True, NAVY),
    ]
    add_multitext(slide, Inches(0.85), Inches(2.7), Inches(5.8), Inches(4.2),
                  lines=principle_lines, line_spacing=1.65)

    # 우측: Part 구조 표
    add_textbox(slide, Inches(7.2), Inches(2.05), Inches(5.5), Inches(0.4),
                text="Part 0 ~ Part 7 구조표", size=12, bold=True, color=NAVY)
    add_rect(slide, Inches(7.2), Inches(2.45), Inches(0.8), Inches(0.03),
             fill_color=GOLD, line_color=None)

    parts = [
        ("Part 0",   "10p",  "리포트 발췌 (손글씨 7항목, 4SE 2p 분할)"),
        ("Part 1",   "12p",  "13영역 인생 지도 (압축)"),
        ("연간",     "2p",   "1년 한눈에 (시장 표준)"),
        ("Part 2",   "10p",  "연간 비전 · 90일 마일스톤"),
        ("월간",     "24p",  "12개월 × 2p (시장 표준)"),
        ("Part 3",   "104p", "주간 펼침면 52주 × 2p (차별점)"),
        ("Part 4",   "26p",  "영역별 분기 회고 (압축)"),
        ("Part 5",   "12p",  "감사 일기 (월별)"),
        ("Part 6",   "8p",   "1:1 코칭 약속 · 진척 기록"),
        ("Part 7",   "24p",  "부록 7p + Owner Profile 1p + 자유 메모장 16p (일일 대체)"),
        ("여유 간지","16p",  "파트 간지 + 16절판 균형용"),
    ]

    # 표 헤더
    th_y = Inches(2.7)
    add_rect(slide, Inches(7.2), th_y, Inches(5.5), Inches(0.4),
             fill_color=NAVY, line_color=None)
    add_textbox(slide, Inches(7.3), th_y + Inches(0.07), Inches(1.0), Inches(0.3),
                text="Part",  size=10, bold=True, color=WHITE)
    add_textbox(slide, Inches(8.4), th_y + Inches(0.07), Inches(0.9), Inches(0.3),
                text="페이지", size=10, bold=True, color=WHITE)
    add_textbox(slide, Inches(9.4), th_y + Inches(0.07), Inches(3.2), Inches(0.3),
                text="내용",   size=10, bold=True, color=WHITE)

    # 표 본문 (v1.3: 10행이므로 row_h 축소)
    row_h = Inches(0.34)
    for i, (p, pg, desc) in enumerate(parts):
        y = th_y + Inches(0.4) + row_h * i
        bg = GRAY_50 if i % 2 == 0 else WHITE
        add_rect(slide, Inches(7.2), y, Inches(5.5), row_h,
                 fill_color=bg, line_color=GRAY_200, line_width=0.4)
        add_textbox(slide, Inches(7.3), y + Inches(0.05), Inches(1.0), Inches(0.28),
                    text=p, size=9.5, bold=True, color=NAVY)
        add_textbox(slide, Inches(8.4), y + Inches(0.05), Inches(0.9), Inches(0.28),
                    text=pg, size=9.5, color=GRAY_700)
        add_textbox(slide, Inches(9.4), y + Inches(0.05), Inches(3.2), Inches(0.28),
                    text=desc, size=9, color=GRAY_900)

    # 합계 행
    sum_y = th_y + Inches(0.4) + row_h * len(parts)
    add_rect(slide, Inches(7.2), sum_y, Inches(5.5), Inches(0.36),
             fill_color=CREAM, line_color=None)
    add_textbox(slide, Inches(7.3), sum_y + Inches(0.06), Inches(1.0), Inches(0.28),
                text="합계", size=10, bold=True, color=NAVY)
    add_textbox(slide, Inches(8.4), sum_y + Inches(0.06), Inches(2.2), Inches(0.28),
                text="240p + 16p = 256p", size=10, bold=True, color=NAVY)
    add_textbox(slide, Inches(10.6), sum_y + Inches(0.06), Inches(2.1), Inches(0.28),
                text="16의 배수 · 16절판", size=9.5, color=GRAY_700)

    add_footer(slide, 8)


# ============================================================
# 슬라이드 9: Part 0 — VDP
# ============================================================
def slide_09_part0(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 9, "03 · 속지 구조", "Part 0 — 리포트 발췌 (손글씨 7항목 · 10p · 4SE 2p 분할 v1.3)")

    # 상단 배지
    add_rect(slide, Inches(0.85), Inches(2.05), Inches(2.8), Inches(0.45),
             fill_color=GOLD, line_color=None)
    add_textbox(slide, Inches(0.85), Inches(2.13), Inches(2.8), Inches(0.3),
                text="100% 자기화 — 손글씨", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 핵심 설명
    add_textbox(slide, Inches(3.9), Inches(2.15), Inches(8.8), Inches(0.35),
                text="사용자가 리포트 결과를 손으로 자기 다이어리에 옮겨 적는 의식 — 자기 관련 기억 +α (Mueller & Oppenheimer, 2014)",
                size=10.5, color=GRAY_700)

    # 좌측: 인쇄될 항목 7가지
    add_textbox(slide, Inches(0.85), Inches(2.85), Inches(6.0), Inches(0.4),
                text="빈칸 라벨 + 손글씨 영역 (7항목 × 1p)", size=12, bold=True, color=NAVY)
    add_rect(slide, Inches(0.85), Inches(3.25), Inches(1.0), Inches(0.03),
             fill_color=GOLD, line_color=None)

    items = [
        ("1.  사명 (Mission)",        "핵심 문장 1줄 + 보조 문장 1줄"),
        ("2.  비전 (Vision)",         "핵심 문장 1줄 + 보조 문장 1줄"),
        ("3.  4SE 응답 강도(%) 2p",   "Page A: 자기이해+자기표현 / Page B: 자기설계+자기실행 (v1.3)"),
        ("4.  TOP3 강점",             "성장 가이드맵에서 발췌"),
        ("5.  TOP2 성장 포인트",      "성장 가이드맵에서 발췌"),
        ("6.  실행 프로파일 6필드",    "유형/스타일/추진력/몰입환경/활동/도구"),
        ("7.  추천 진로 3카드",       "추천 진로 · 추천 교육 · 확장 방향"),
    ]
    y0 = Inches(3.45)
    row_h = Inches(0.42)
    for i, (k, v) in enumerate(items):
        y = y0 + row_h * i
        add_textbox(slide, Inches(0.85), y, Inches(2.8), Inches(0.4),
                    text=k, size=10.5, bold=True, color=GRAY_900)
        add_textbox(slide, Inches(3.7), y, Inches(3.5), Inches(0.4),
                    text=v, size=10, color=GRAY_500)

    # 우측: 견적 요청 안내
    add_rect(slide, Inches(7.6), Inches(2.85), Inches(5.13), Inches(3.85),
             fill_color=GRAY_50, line_color=GRAY_200, line_width=0.5)
    add_rect(slide, Inches(7.6), Inches(2.85), Inches(0.06), Inches(3.85),
             fill_color=GOLD, line_color=None)

    note_lines = [
        ("제조사 견적 요청 사항", 12, True, NAVY),
        ("", 6, False, WHITE),
        ("· 70g 종이 — 손글씨 친화도 테스트", 10.5, False, GRAY_900),
        ("  (만년필·중성펜 번짐 없음)", 10.5, False, GRAY_900),
        ("· 펼침 180° 평평도 (양장 사철) — 양면 필기 안정", 10.5, False, GRAY_900),
        ("· 빈칸 행 간격 7mm — 손글씨 가독성", 10.5, False, GRAY_900),
        ("· 라벨 서체: Noto Sans CJK KR Light 9pt", 10.5, False, GRAY_900),
        ("· 단일 일반 인쇄 — VDP 없음, 전수량 동일", 10.5, False, GRAY_900),
        ("", 6, False, WHITE),
        ("→ 사양서 PDF v1.3 §4 와 동일", 10, True, GOLD),
    ]
    add_multitext(slide, Inches(7.85), Inches(2.95), Inches(4.85), Inches(3.7),
                  lines=note_lines, line_spacing=1.5)

    add_footer(slide, 9)


# ============================================================
# 슬라이드 10: Part 1·2 — 13영역 + 연간비전
# ============================================================
def slide_10_part12(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 10, "03 · 속지 구조", "Part 1 · 2 — 13영역 인생지도 + 연간 비전 (28p)")

    # 좌측: Part 1
    add_rect(slide, Inches(0.6), Inches(2.05), Inches(6.0), Inches(4.85),
             fill_color=WHITE, line_color=GRAY_200, line_width=1.0)
    add_rect(slide, Inches(0.6), Inches(2.05), Inches(6.0), Inches(0.6),
             fill_color=NAVY, line_color=None)
    add_textbox(slide, Inches(0.85), Inches(2.17), Inches(5.7), Inches(0.4),
                text="Part 1 — 13영역 인생 지도 (16p)", size=13, bold=True, color=WHITE)

    p1_lines = [
        ("페이지 구성", 11, True, NAVY),
        ("· 13영역 명칭만 인쇄 (인생포트폴리오 IP)", 10.5, False, GRAY_900),
        ("· 점수 칸 / 메모 칸은 빈칸", 10.5, False, GRAY_900),
        ("· 1영역당 약 1.2p 할당", 10.5, False, GRAY_900),
        ("", 6, False, WHITE),
        ("13영역 라벨 (보고서 6.2 그대로)", 11, True, NAVY),
        ("신앙·소명 / 가족·관계 / 직업·경력", 10, False, GRAY_900),
        ("재정·자산 / 건강·체력 / 학습·성장", 10, False, GRAY_900),
        ("취미·여가 / 봉사·기여 / 자기관리", 10, False, GRAY_900),
        ("거주·환경 / 인간관계 / 정서·내면 / 유산·미래", 10, False, GRAY_900),
    ]
    add_multitext(slide, Inches(0.85), Inches(2.85), Inches(5.5), Inches(3.9),
                  lines=p1_lines, line_spacing=1.55)

    # 우측: Part 2
    add_rect(slide, Inches(6.8), Inches(2.05), Inches(6.0), Inches(4.85),
             fill_color=WHITE, line_color=GRAY_200, line_width=1.0)
    add_rect(slide, Inches(6.8), Inches(2.05), Inches(6.0), Inches(0.6),
             fill_color=NAVY, line_color=None)
    add_textbox(slide, Inches(7.05), Inches(2.17), Inches(5.7), Inches(0.4),
                text="Part 2 — 연간 비전 · 90일 마일스톤 (12p)", size=13, bold=True, color=WHITE)

    p2_lines = [
        ("페이지 구성", 11, True, NAVY),
        ("· “VISION” 박스 + “분기 마일스톤 3개” 라벨만 인쇄", 10.5, False, GRAY_900),
        ("· 빈칸은 사용자가 검사 리포트 비전 카드 보고 옮겨 적음", 10.5, False, GRAY_900),
        ("· Hyatt 스타일 — 연간 → 분기 → 90일 분할", 10.5, False, GRAY_900),
        ("", 6, False, WHITE),
        ("매핑되는 검사 리포트 항목", 11, True, NAVY),
        ("· 리포트 페이지 1 — 비전 카드의 핵심·보조 문장", 10, False, GRAY_900),
        ("· 실행 프로그램 페이지 1 — 분기 테마 카드 1개", 10, False, GRAY_900),
        ("· 실행 프로그램 페이지 3 — 3개월 카드 3개", 10, False, GRAY_900),
        ("", 6, False, WHITE),
        ("→ 라벨·항목명 운영 PDF와 100% 정합", 10.5, True, GOLD),
    ]
    add_multitext(slide, Inches(7.05), Inches(2.85), Inches(5.5), Inches(3.9),
                  lines=p2_lines, line_spacing=1.55)

    add_footer(slide, 10)


# ============================================================
# 슬라이드 11: Part 3 — 주간 펼침면
# ============================================================
def slide_11_part3(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 11, "03 · 속지 구조", "Part 3 — 주간 펼침면 52주 × 2p (104p)")

    # 상단 메시지
    add_textbox(slide, Inches(0.85), Inches(2.05), Inches(12.0), Inches(0.4),
                text="몰스킨 위클리 + 프랭클린 우선순위 하이브리드 — 좌측 계획 / 우측 메모 (자체 IP 5대 차별화 적용)",
                size=11, color=GRAY_700)

    # 좌우 펼침면 mock-up
    lx = Inches(0.85)
    rx = Inches(6.95)
    py = Inches(2.6)
    pw = Inches(5.9)
    ph = Inches(3.95)

    # ----- 좌측 페이지 — 계획 -----
    add_rect(slide, lx, py, pw, ph,
             fill_color=GRAY_50, line_color=GRAY_200, line_width=0.7)
    add_textbox(slide, lx + Inches(0.2), py + Inches(0.08), pw - Inches(0.4), Inches(0.3),
                text="LEFT — 이번 주 우선순위 · 일정", size=9, bold=True, color=GOLD)
    add_textbox(slide, lx + Inches(0.2), py + Inches(0.34), pw - Inches(0.4), Inches(0.4),
                text="Week ##", size=14, bold=True, color=NAVY)

    # 이번 주 사명 한 줄 (자체 IP ①)
    add_textbox(slide, lx + Inches(0.2), py + Inches(0.75), Inches(3.0), Inches(0.25),
                text="① 이번 주 사명 (1줄)", size=8.5, bold=True, color=GOLD)
    add_rect(slide, lx + Inches(0.2), py + Inches(1.02), pw - Inches(0.4), Emu(7000),
             fill_color=NAVY, line_color=None)

    # A/B/C 우선순위 (프랭클린 참고)
    add_textbox(slide, lx + Inches(0.2), py + Inches(1.15), Inches(4.5), Inches(0.25),
                text="② A · B · C 우선순위 (프랭클린 참고)", size=8.5, bold=True, color=GOLD)
    labels = ["A", "B", "C"]
    for i, lab in enumerate(labels):
        ly = py + Inches(1.42) + Inches(0.28) * i
        # 라벨 박스
        add_rect(slide, lx + Inches(0.2), ly, Inches(0.28), Inches(0.22),
                 fill_color=NAVY, line_color=None)
        add_textbox(slide, lx + Inches(0.2), ly - Inches(0.01), Inches(0.28), Inches(0.22),
                    text=lab, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # 빈칸
        add_rect(slide, lx + Inches(0.55), ly + Inches(0.18),
                 pw - Inches(0.75), Emu(3500),
                 fill_color=GRAY_200, line_color=None)

    # If-Then 3개 (자체 IP ③ — Gollwitzer)
    add_textbox(slide, lx + Inches(0.2), py + Inches(2.35), Inches(5.0), Inches(0.25),
                text="③ If-Then 3개 (Gollwitzer 실행 의도)", size=8.5, bold=True, color=GOLD)
    for i in range(3):
        ly = py + Inches(2.62) + Inches(0.22) * i
        add_rect(slide, lx + Inches(0.4), ly + Inches(0.13),
                 pw - Inches(0.6), Emu(3500),
                 fill_color=GRAY_200, line_color=None)

    # 요일별 일정 (몰스킨 참고)
    add_textbox(slide, lx + Inches(0.2), py + Inches(3.32), Inches(5.0), Inches(0.25),
                text="④ 요일별 일정 (월→일 간소)", size=8.5, bold=True, color=GOLD)
    days_short = "월 · 화 · 수 · 목 · 금 · 토 · 일"
    add_textbox(slide, lx + Inches(0.2), py + Inches(3.55), pw - Inches(0.4), Inches(0.3),
                text=days_short, size=9, color=GRAY_700)

    # ----- 우측 페이지 — 메모 + 회고 -----
    add_rect(slide, rx, py, pw, ph,
             fill_color=GRAY_50, line_color=GRAY_200, line_width=0.7)
    add_textbox(slide, rx + Inches(0.2), py + Inches(0.08), pw - Inches(0.4), Inches(0.3),
                text="RIGHT — 자유 메모 + 이번 주 회고", size=9, bold=True, color=GOLD)
    add_textbox(slide, rx + Inches(0.2), py + Inches(0.34), pw - Inches(0.4), Inches(0.4),
                text="도트 그리드 7mm + 회고 3줄", size=14, bold=True, color=NAVY)

    # 도트 그리드 영역 (몰스킨 참고)
    add_textbox(slide, rx + Inches(0.2), py + Inches(0.85), Inches(5.0), Inches(0.25),
                text="· 7mm 도트 그리드 자유 메모 영역", size=8.5, bold=True, color=GOLD)
    # 도트 시각화 (점 계열)
    dot_area_y = py + Inches(1.15)
    dot_area_h = Inches(1.75)
    add_rect(slide, rx + Inches(0.2), dot_area_y, pw - Inches(0.4), dot_area_h,
             fill_color=WHITE, line_color=GRAY_200, line_width=0.4)
    # 도트 파턴 (8행 × 18열)
    dot_color = RGBColor(0xC0, 0xC4, 0xC8)
    cols_n = 18
    rows_n = 8
    for ri in range(rows_n):
        for ci in range(cols_n):
            cx = rx + Inches(0.4) + Inches(0.28) * ci
            cy = dot_area_y + Inches(0.18) + Inches(0.2) * ri
            add_rect(slide, cx, cy, Emu(20000), Emu(20000),
                     fill_color=dot_color, line_color=None)

    # 이번 주 회고 3줄 (자체 IP ⑤)
    add_textbox(slide, rx + Inches(0.2), py + Inches(3.05), Inches(5.0), Inches(0.25),
                text="⑤ 이번 주 회고 3줄 (자체 IP)", size=8.5, bold=True, color=GOLD)
    review_labels = ["잘된 것", "배운 것", "다음 주"]
    for i, lab in enumerate(review_labels):
        ly = py + Inches(3.3) + Inches(0.22) * i
        add_textbox(slide, rx + Inches(0.2), ly, Inches(0.85), Inches(0.22),
                    text=lab, size=8.5, bold=True, color=NAVY)
        add_rect(slide, rx + Inches(1.1), ly + Inches(0.15),
                 pw - Inches(1.3), Emu(3500),
                 fill_color=GRAY_200, line_color=None)

    # 하단 차별화 메시지
    add_textbox(slide, Inches(0.85), Inches(6.7), Inches(12.0), Inches(0.3),
                text="자체 IP 차별화: ① 주간 사명 · ② 프랭클린 우선순위 · ③ Gollwitzer If-Then · ④ 몰스킨 서셋 · ⑤ 자체 3줄 회고",
                size=9, color=GRAY_500)

    add_footer(slide, 11)


# ============================================================
# 슬라이드 12: Part 4·5 — 분기 회고 + 감사 일기
# ============================================================
def slide_12_part45(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 12, "03 · 속지 구조", "Part 4 · 5 — 분기 회고 + 감사 일기 (64p)")

    # 좌: Part 4
    add_rect(slide, Inches(0.6), Inches(2.05), Inches(6.0), Inches(4.85),
             fill_color=WHITE, line_color=GRAY_200, line_width=1.0)
    add_rect(slide, Inches(0.6), Inches(2.05), Inches(6.0), Inches(0.6),
             fill_color=NAVY, line_color=None)
    add_textbox(slide, Inches(0.85), Inches(2.17), Inches(5.7), Inches(0.4),
                text="Part 4 — 영역별 분기 회고 (4p × 13 = 52p)", size=13, bold=True, color=WHITE)

    p4_lines = [
        ("구성", 11, True, NAVY),
        ("· 영역명 인쇄(인생포트폴리오 13영역)", 10.5, False, GRAY_900),
        ("· Pennebaker 표현적 글쓰기 4문항 빈칸", 10.5, False, GRAY_900),
        ("· 1영역당 4p (분기당 1p 회고)", 10.5, False, GRAY_900),
        ("", 6, False, WHITE),
        ("Pennebaker 4문항(고정 라벨)", 11, True, NAVY),
        ("Q1. 지난 분기 이 영역에서 일어난 사건", 10, False, GRAY_900),
        ("Q2. 그때 느낀 감정과 생각", 10, False, GRAY_900),
        ("Q3. 그것이 내게 의미하는 바", 10, False, GRAY_900),
        ("Q4. 다음 분기의 의도와 행동", 10, False, GRAY_900),
        ("", 6, False, WHITE),
        ("근거: Pennebaker(1986) — Expressive Writing", 9.5, False, GRAY_500),
    ]
    add_multitext(slide, Inches(0.85), Inches(2.85), Inches(5.5), Inches(3.9),
                  lines=p4_lines, line_spacing=1.5)

    # 우: Part 5
    add_rect(slide, Inches(6.8), Inches(2.05), Inches(6.0), Inches(4.85),
             fill_color=WHITE, line_color=GRAY_200, line_width=1.0)
    add_rect(slide, Inches(6.8), Inches(2.05), Inches(6.0), Inches(0.6),
             fill_color=NAVY, line_color=None)
    add_textbox(slide, Inches(7.05), Inches(2.17), Inches(5.7), Inches(0.4),
                text="Part 5 — 감사 일기 월별 (12p)", size=13, bold=True, color=WHITE)

    p5_lines = [
        ("구성", 11, True, NAVY),
        ("· 월 1p (12개월)", 10.5, False, GRAY_900),
        ("· 라벨 \"이번 달 감사 3가지\" 인쇄", 10.5, False, GRAY_900),
        ("· 빈칸 3줄 + 우측 짧은 메모 칸", 10.5, False, GRAY_900),
        ("", 6, False, WHITE),
        ("디자인 원칙", 11, True, NAVY),
        ("· 페이지 상단에 월 / 년 표기만 인쇄", 10, False, GRAY_900),
        ("· 빈 공간 비율 70% 이상 — 사용자 공간 확보", 10, False, GRAY_900),
        ("· 골드 라인 1줄로 시각적 분리", 10, False, GRAY_900),
        ("", 6, False, WHITE),
        ("근거: Cregg & Cheavens(2024) — Gratitude RCT", 9.5, False, GRAY_500),
    ]
    add_multitext(slide, Inches(7.05), Inches(2.85), Inches(5.5), Inches(3.9),
                  lines=p5_lines, line_spacing=1.5)

    add_footer(slide, 12)


# ============================================================
# 슬라이드 13: Part 6·7 — 코칭기록 + 부록
# ============================================================
def slide_13_part67(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 13, "03 · 속지 구조", "Part 6 · 7 — 1:1 코칭 기록 + 부록 (20p)")

    # 좌: Part 6 (성과 추적 보드 표)
    add_rect(slide, Inches(0.6), Inches(2.05), Inches(6.0), Inches(4.85),
             fill_color=WHITE, line_color=GRAY_200, line_width=1.0)
    add_rect(slide, Inches(0.6), Inches(2.05), Inches(6.0), Inches(0.6),
             fill_color=NAVY, line_color=None)
    add_textbox(slide, Inches(0.85), Inches(2.17), Inches(5.7), Inches(0.4),
                text="Part 6 — 1:1 코칭 약속 · 진척 기록 (8p)", size=13, bold=True, color=WHITE)

    add_textbox(slide, Inches(0.85), Inches(2.85), Inches(5.5), Inches(0.4),
                text="성과 추적 보드 — 4컬럼 표 (실행 프로그램 PDF p4 동일)",
                size=10.5, bold=True, color=NAVY)

    # 4컬럼 표 mock
    cols = [("주차", 0.7), ("실행 과제", 2.2), ("완료(Y/N)", 1.0), ("성찰 메모", 1.6)]
    th_y = Inches(3.3)
    th_x = Inches(0.85)

    cx = th_x
    for name, w in cols:
        add_rect(slide, cx, th_y, Inches(w), Inches(0.4),
                 fill_color=NAVY, line_color=None)
        add_textbox(slide, cx, th_y + Inches(0.08), Inches(w), Inches(0.3),
                    text=name, size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += Inches(w)

    # 빈 행 6개
    for i in range(6):
        ry = th_y + Inches(0.4) + Inches(0.35) * i
        cx = th_x
        bg = GRAY_50 if i % 2 == 0 else WHITE
        for name, w in cols:
            add_rect(slide, cx, ry, Inches(w), Inches(0.35),
                     fill_color=bg, line_color=GRAY_200, line_width=0.4)
            cx += Inches(w)

    add_textbox(slide, Inches(0.85), Inches(6.0), Inches(5.7), Inches(0.6),
                text="근거: Matthews(2007) Group 5 — 글로 쓴 목표 + 친구에게 약속·진척 보고 시 달성률 +78%",
                size=9, color=GRAY_500, line_spacing=1.4)

    # 우: Part 7 (부록)
    add_rect(slide, Inches(6.8), Inches(2.05), Inches(6.0), Inches(4.85),
             fill_color=WHITE, line_color=GRAY_200, line_width=1.0)
    add_rect(slide, Inches(6.8), Inches(2.05), Inches(6.0), Inches(0.6),
             fill_color=NAVY, line_color=None)
    add_textbox(slide, Inches(7.05), Inches(2.17), Inches(5.7), Inches(0.4),
                text="Part 7 — 부록 (12p)", size=13, bold=True, color=WHITE)

    p7_lines = [
        ("구성 — 라벨만 인쇄", 11, True, NAVY),
        ("· 성경구절 또는 명언 모음 (4p)", 10.5, False, GRAY_900),
        ("· 13영역 참고 가이드 (4p)", 10.5, False, GRAY_900),
        ("· 사용 가이드 요약 (2p)", 10.5, False, GRAY_900),
        ("· 연락처 / 라이선스 / 색인 (2p)", 10.5, False, GRAY_900),
        ("", 6, False, WHITE),
        ("신학 안전 표현 원칙", 11, True, NAVY),
        ("· \"성경구절 또는 명언 선택 가능\" 중립 표현", 10, False, GRAY_900),
        ("· Frankl·Dik·Wrzesniewski 등 학술 인용은", 10, False, GRAY_900),
        ("  Pre-evangelism / Common Grace 영역에서 안전", 10, False, GRAY_900),
        ("· 자기 너머의 부름을 명시적 vocation으로 동일시 않음", 10, False, GRAY_900),
    ]
    add_multitext(slide, Inches(7.05), Inches(2.85), Inches(5.5), Inches(3.9),
                  lines=p7_lines, line_spacing=1.5)

    add_footer(slide, 13)


# ============================================================
# 슬라이드 14: 인쇄 구조 — 2층
# ============================================================
def slide_14_print_structure(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 14, "04 · 인쇄 구조", "인쇄 구조 — 단일 일반 인쇄 (Single-Layer Print)")

    # 상단 메시지
    add_textbox(slide, Inches(0.85), Inches(2.05), Inches(12.0), Inches(0.4),
                text="다이어리 1권 = 전 페이지 동일 일반 인쇄 (속표지·간지 8p + 내지 224p). 구매자별 개인화는 사용자가 손글씨로 수행.",
                size=11, color=GRAY_700)

    # 메인 카드 — 단일 인쇄
    l1_y = Inches(2.7)
    add_rect(slide, Inches(0.6), l1_y, Inches(12.13), Inches(2.0),
             fill_color=GRAY_50, line_color=None)
    add_rect(slide, Inches(0.6), l1_y, Inches(0.1), Inches(2.0),
             fill_color=NAVY, line_color=None)

    add_textbox(slide, Inches(0.9), l1_y + Inches(0.18), Inches(3.0), Inches(0.4),
                text="SINGLE LAYER", size=10, bold=True, color=NAVY)
    add_textbox(slide, Inches(0.9), l1_y + Inches(0.5), Inches(10.5), Inches(0.5),
                text="일반 인쇄 — 빈칸 양식 (라벨만 인쇄 / 사용자가 손글씨로 채움)",
                size=16, bold=True, color=NAVY)

    l1_text = [
        ("· 적용 범위: Part 0 (8p) — 7항목 빈칸 + Part 1~7 (224p) 전부 동일", 10.5, False, GRAY_900),
        ("· 장점: 단일 인쇄 = 단가 안정, 제조사 교섭 단순화, VDP 시스템 투자 불필요", 10.5, False, GRAY_900),
        ("· 견적 요청: 70g 종이 · 256p · 오프셋/디지털 택1", 10.5, True, NAVY),
    ]
    add_multitext(slide, Inches(0.9), l1_y + Inches(1.05),
                  Inches(11.5), Inches(0.85),
                  lines=l1_text, line_spacing=1.5)

    # 하단 — 개인화 메커니즘 (차별화 메시지)
    l2_y = Inches(5.0)
    add_rect(slide, Inches(0.6), l2_y, Inches(12.13), Inches(1.85),
             fill_color=CREAM, line_color=None)
    add_rect(slide, Inches(0.6), l2_y, Inches(0.1), Inches(1.85),
             fill_color=GOLD, line_color=None)

    add_textbox(slide, Inches(0.9), l2_y + Inches(0.15), Inches(8.0), Inches(0.4),
                text="개인화 메커니즘 — 자기화(自己化) by 손글씨", size=14, bold=True, color=NAVY)

    l2_text = [
        ("· Part 0 7항목 (사명/비전/4SE/TOP3/TOP2/실행프로파일/추천진로) — 리포트를 보며 손으로 옮겨 적음", 10.5, False, GRAY_900),
        ("· Part 1~7 전보 — 라벨만 인쇄되어 있고 활용 온도는 사용자의 손글씨에서 창출", 10.5, False, GRAY_900),
        ("· 근거: Mueller & Oppenheimer(2014) — 손글씨 → 개념 이해·장기 기억 우위 / Smoker(2009) — 손글씨 루틴 → 자기 관련 정보 처리 강화", 10, True, NAVY),
    ]
    add_multitext(slide, Inches(0.9), l2_y + Inches(0.65),
                  Inches(11.5), Inches(1.1),
                  lines=l2_text, line_spacing=1.5)

    add_footer(slide, 14)


# ============================================================
# 슬라이드 15: 타이포·여백·페이지 설계 원칙
# ============================================================
def slide_15_typography(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 15, "04 · 인쇄 구조", "타이포 · 여백 · 페이지 설계 원칙")

    # 4 원칙 카드 그리드 2×2
    principles = [
        ("01", "미니멀",        "타이포·여백 위주. 사진·일러스트 최소화. 사용자가 빈 공간을 채움."),
        ("02", "만년필 친화",   "70g 도모에리버 대체재 권장. 비침·번짐 최소. 행간 충분히 확보."),
        ("03", "가독성 우선",   "본문 10.5 ~ 11pt 권장. 한글 행간 1.6 ~ 1.7. 노안 사용자 30~50대 가독성 우선."),
        ("04", "빈 공간 = 가치", "각 페이지 빈 공간 비율 60% 이상. 사용자의 글쓰기 공간이 곧 다이어리의 핵심 IP."),
    ]

    card_w = Inches(5.85)
    card_h = Inches(2.15)
    x0 = Inches(0.6)
    y0 = Inches(2.1)
    gap_x = Inches(0.35)
    gap_y = Inches(0.35)

    for i, (num, title, desc) in enumerate(principles):
        col = i % 2
        row = i // 2
        x = x0 + (card_w + gap_x) * col
        y = y0 + (card_h + gap_y) * row

        add_rect(slide, x, y, card_w, card_h,
                 fill_color=WHITE, line_color=GRAY_200, line_width=1.0)

        # 큰 번호
        add_textbox(slide, x + Inches(0.3), y + Inches(0.15), Inches(1.2), Inches(0.9),
                    text=num, size=42, bold=True, color=GOLD)
        # 골드 라인
        add_rect(slide, x + Inches(0.3), y + Inches(1.05), Inches(0.6), Inches(0.04),
                 fill_color=GOLD, line_color=None)

        # 제목
        add_textbox(slide, x + Inches(1.6), y + Inches(0.3), card_w - Inches(1.8), Inches(0.5),
                    text=title, size=18, bold=True, color=NAVY)
        # 설명
        add_textbox(slide, x + Inches(1.6), y + Inches(0.9), card_w - Inches(1.8), Inches(1.1),
                    text=desc, size=10.5, color=GRAY_700, line_spacing=1.6)

    add_footer(slide, 15)


# ============================================================
# 슬라이드 16: 견적 요청 항목 + 회신 양식
# ============================================================
def slide_16_quote(prs):
    slide = add_blank_slide(prs)
    add_slide_header(slide, 16, "05 · 견적 요청", "견적 요청 항목 + 회신 양식")

    # 좌측: 8개 항목 그룹화 (v1.3)
    add_textbox(slide, Inches(0.85), Inches(2.05), Inches(6.0), Inches(0.4),
                text="견적 요청 항목 — 8개 (사양서 v1.3 §4 ~ §6)", size=12, bold=True, color=NAVY)
    add_rect(slide, Inches(0.85), Inches(2.45), Inches(1.0), Inches(0.03),
             fill_color=GOLD, line_color=None)

    groups = [
        ("외관(2)",     "A5 양장 사철 / PU + 금박·형압·가름끈 2 부속 일괄"),
        ("종이(1)",     "70g 만년필 친화 — 단일 시나리오 견적"),
        ("페이지(1)",   "256p (속표지·간지 8p + 본문 248p)"),
        ("인쇄(1)",     "일반 인쇄 1종 — 오프셋/디지털 추천 선택"),
        ("후가공(1)",   "금박(C9A04F) + 형압 로고"),
        ("포장(1)",     "박스 + 인삿카드 + 사용가이드 8p"),
        ("납기(1)",     "샘플 · 본인쇄 일정 + 예약 옵션"),
    ]
    y0 = Inches(2.7)
    row_h = Inches(0.55)
    for i, (k, v) in enumerate(groups):
        y = y0 + row_h * i
        # 라벨
        add_rect(slide, Inches(0.85), y, Inches(1.1), Inches(0.4),
                 fill_color=NAVY, line_color=None)
        add_textbox(slide, Inches(0.85), y + Inches(0.08), Inches(1.1), Inches(0.3),
                    text=k, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # 설명
        add_textbox(slide, Inches(2.05), y + Inches(0.05), Inches(4.8), Inches(0.4),
                    text=v, size=10, color=GRAY_900, line_spacing=1.4)

    # 우측: 회신 양식 안내 (v1.3 — 단일 시나리오)
    rx = Inches(7.2)
    add_textbox(slide, rx, Inches(2.05), Inches(5.5), Inches(0.4),
                text="회신 양식 — 단일 시나리오 견적표", size=12, bold=True, color=NAVY)
    add_rect(slide, rx, Inches(2.45), Inches(1.0), Inches(0.03),
             fill_color=GOLD, line_color=None)

    # 단일 시나리오 표 (v1.3)
    scenarios = [
        ("항목",      "수량",   "견적·납기"),
        ("샘플 제작", "1부",    "본 계약 이전"),
        ("본 제작",  "500부",  "1차 발주"),
        ("추가 제작", "+500부", "재주문 옵션"),
        ("단가",      "1부",    "도매 서적용"),
    ]
    cw = [Inches(1.4), Inches(1.3), Inches(2.8)]
    ty = Inches(2.7)
    for i, row in enumerate(scenarios):
        cx = rx
        bg = NAVY if i == 0 else (GRAY_50 if i % 2 == 0 else WHITE)
        tc = WHITE if i == 0 else GRAY_900
        bold = (i == 0)
        for j, cell in enumerate(row):
            add_rect(slide, cx, ty, cw[j], Inches(0.42),
                     fill_color=bg, line_color=GRAY_200, line_width=0.4)
            add_textbox(slide, cx + Inches(0.1), ty + Inches(0.1),
                        cw[j] - Inches(0.2), Inches(0.3),
                        text=cell, size=10, bold=bold, color=tc)
            cx += cw[j]
        ty += Inches(0.42)

    # 회신 일정 박스
    add_rect(slide, rx, Inches(5.05), Inches(5.5), Inches(1.65),
             fill_color=CREAM, line_color=None)
    add_rect(slide, rx, Inches(5.05), Inches(0.06), Inches(1.65),
             fill_color=GOLD, line_color=None)

    sched_lines = [
        ("회신 요청 일정", 11, True, NAVY),
        ("· 사양서 수령 후 7영업일 이내 회신", 10.5, False, GRAY_900),
        ("· 4 시나리오 견적 + VDP 가능 여부 + 납기일", 10.5, False, GRAY_900),
        ("· 회신처: [추후 협의] / 담당자: [추후 협의]", 10.5, False, GRAY_700),
        ("· 비교 후 1개 제조사 샘플 제작 의뢰 진행", 10.5, True, NAVY),
    ]
    add_multitext(slide, rx + Inches(0.2), Inches(5.15),
                  Inches(5.2), Inches(1.5),
                  lines=sched_lines, line_spacing=1.5)

    add_footer(slide, 16)


# ============================================================
# 메인
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 16 슬라이드 빌드
    slide_01_cover(prs)
    slide_02_toc(prs)
    slide_03_definition(prs)
    slide_04_target(prs)
    slide_05_appearance(prs)
    slide_06_color(prs)
    slide_07_paper(prs)
    slide_08_inner_overview(prs)
    slide_09_part0(prs)
    slide_10_part12(prs)
    slide_11_part3(prs)
    slide_12_part45(prs)
    slide_13_part67(prs)
    slide_14_print_structure(prs)
    slide_15_typography(prs)
    slide_16_quote(prs)

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))

    size_kb = OUTPUT_PPTX.stat().st_size / 1024
    print(f"[OK] PPT 생성 완료: {OUTPUT_PPTX}")
    print(f"     크기: {size_kb:.1f} KB")
    print(f"     슬라이드: {len(prs.slides)}장")


if __name__ == "__main__":
    main()
